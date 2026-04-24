"""Core PPTX assembly: create_pptx_with_native_svg."""

from __future__ import annotations

import hashlib
import re
import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from .drawingml_converter import convert_svg_to_slide_shapes
from .pptx_dimensions import (
    CANVAS_FORMATS,
    get_slide_dimensions, get_pixel_dimensions,
    get_viewbox_dimensions, detect_format_from_svg,
)
from .pptx_media import (
    PNG_RENDERER,
    get_png_renderer_info, convert_svg_to_png,
)
from .pptx_notes import (
    markdown_to_plain_text,
    create_notes_slide_xml, create_notes_slide_rels_xml,
)
from .pptx_slide_xml import (
    ANIMATIONS_AVAILABLE, TRANSITIONS,
    create_slide_xml_with_svg, create_slide_rels_xml,
)

# Re-import create_transition_xml only if available
try:
    from pptx_animations import create_transition_xml
except ImportError:
    create_transition_xml = None


def _append_relationship(
    rels_path: Path,
    rel_type: str,
    target: str,
) -> str:
    """Append a relationship entry with the next available rId."""
    with open(rels_path, 'r', encoding='utf-8') as f:
        rels_content = f.read()

    rid_numbers = [int(match) for match in re.findall(r'Id="rId(\d+)"', rels_content)]
    next_rid = f'rId{max(rid_numbers, default=0) + 1}'
    rel_xml = (
        f'  <Relationship Id="{next_rid}" '
        f'Type="{rel_type}" Target="{target}"/>'
    )
    rels_content = rels_content.replace(
        '</Relationships>', rel_xml + '\n</Relationships>',
    )

    with open(rels_path, 'w', encoding='utf-8') as f:
        f.write(rels_content)

    return next_rid


def create_pptx_with_native_svg(
    svg_files: list[Path],
    output_path: Path,
    canvas_format: str | None = None,
    verbose: bool = True,
    transition: str | None = 'fade',
    transition_duration: float = 0.5,
    auto_advance: float | None = None,
    use_compat_mode: bool = True,
    notes: dict[str, str] | None = None,
    enable_notes: bool = True,
    use_native_shapes: bool = False,
) -> bool:
    """Create a PPTX file with native SVG.

    Args:
        svg_files: List of SVG files.
        output_path: Output PPTX path.
        canvas_format: Canvas format key.
        verbose: Whether to output detailed information.
        transition: Transition effect name.
        transition_duration: Transition duration in seconds.
        auto_advance: Auto-advance interval in seconds.
        use_compat_mode: Use Office compatibility mode (PNG + SVG dual format).
        notes: Notes dict, key is SVG stem, value is notes content.
        enable_notes: Whether to enable notes embedding.
        use_native_shapes: Convert SVG to native DrawingML shapes.

    Returns:
        Whether all slides were successfully created.
    """
    if not svg_files:
        print("Error: No SVG files found")
        return False

    # Native shapes mode takes priority over compat mode
    if use_native_shapes:
        use_compat_mode = False

    # Check compatibility mode dependencies
    renderer_name, renderer_status, renderer_hint = get_png_renderer_info()
    if not use_native_shapes and use_compat_mode and PNG_RENDERER is None:
        print("Warning: No PNG rendering library installed, cannot use compatibility mode")
        print(f"  {renderer_hint}")
        print("  Will use pure SVG mode (may not display in Office LTSC 2021 and similar versions)")
        use_compat_mode = False

    # Auto-detect canvas format or get dimensions from viewBox
    custom_pixels: tuple[int, int] | None = None
    if canvas_format is None:
        canvas_format = detect_format_from_svg(svg_files[0])
        if canvas_format and verbose:
            format_name = CANVAS_FORMATS.get(canvas_format, {}).get('name', canvas_format)
            print(f"  Detected canvas format: {format_name}")

    if canvas_format is None:
        custom_pixels = get_viewbox_dimensions(svg_files[0])
        if custom_pixels and verbose:
            print(f"  Using SVG viewBox dimensions: {custom_pixels[0]} x {custom_pixels[1]} px")

    if canvas_format is None and custom_pixels is None:
        canvas_format = 'ppt169'
        if verbose:
            print(f"  Using default format: PPT 16:9")

    width_emu, height_emu = get_slide_dimensions(canvas_format or 'ppt169', custom_pixels)
    pixel_width, pixel_height = get_pixel_dimensions(canvas_format or 'ppt169', custom_pixels)

    if verbose:
        print(f"  Slide dimensions: {pixel_width} x {pixel_height} px")
        print(f"  SVG file count: {len(svg_files)}")
        if use_native_shapes:
            print(f"  Mode: Native DrawingML shapes (directly editable)")
        elif use_compat_mode:
            print(f"  Compatibility mode: Enabled (PNG + SVG dual format)")
            print(f"  PNG renderer: {renderer_name} {renderer_status}")
        else:
            print(f"  Compatibility mode: Disabled (pure SVG)")
        if transition:
            trans_name = TRANSITIONS.get(transition, {}).get('name', transition) if TRANSITIONS else transition
            print(f"  Transition effect: {trans_name}")
        if enable_notes and notes:
            print(f"  Speaker notes: {len(notes)} page(s)")
        elif enable_notes:
            print(f"  Speaker notes: Enabled (no notes files found)")
        else:
            print(f"  Speaker notes: Disabled")
        print()

    temp_dir = Path(tempfile.mkdtemp())

    try:
        # Create base PPTX with python-pptx
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu

        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)

        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))

        # Extract PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)

        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)

        success_count = 0
        has_any_image = False
        media_cache: dict[tuple[str, str], str] = {}
        notes_slides_created: set[int] = set()

        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i

            try:
                # ---- Native shapes mode ----
                if use_native_shapes:
                    slide_xml, media_files_dict, rel_entries = convert_svg_to_slide_shapes(
                        svg_path, slide_num=slide_num, verbose=verbose,
                    )

                    # Add transition if specified
                    if transition and ANIMATIONS_AVAILABLE and create_transition_xml:
                        transition_xml = '\n' + create_transition_xml(
                            effect=transition,
                            duration=transition_duration,
                            advance_after=auto_advance,
                        )
                        slide_xml = slide_xml.replace(
                            '</p:sld>',
                            transition_xml + '\n</p:sld>',
                        )

                    # Write slide XML
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    # Write media files
                    media_name_map: dict[str, str] = {}
                    for media_name, media_data in media_files_dict.items():
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        media_hash = hashlib.sha256(media_data).hexdigest()
                        cache_key = (ext, media_hash)
                        cached_name = media_cache.get(cache_key)

                        if cached_name is None:
                            cached_name = media_name
                            media_cache[cache_key] = cached_name
                            with open(media_dir / cached_name, 'wb') as f:
                                f.write(media_data)

                        media_name_map[media_name] = cached_name

                    for rel in rel_entries:
                        target = rel.get('target', '')
                        if not target.startswith('../media/'):
                            continue
                        media_name = target.split('../media/', 1)[1]
                        mapped_name = media_name_map.get(media_name)
                        if mapped_name:
                            rel['target'] = f'../media/{mapped_name}'

                    # Build relationships XML
                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'

                    extra_rels = ''
                    for rel in rel_entries:
                        extra_rels += (
                            f'\n  <Relationship Id="{rel["id"]}" '
                            f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
                        )

                    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>{extra_rels}
</Relationships>'''
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                    # Track image formats for Content_Types
                    for media_name in media_files_dict:
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        if ext in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
                            has_any_image = True

                # ---- Legacy SVG embedding mode ----
                else:
                    svg_filename = f'image{i}.svg'
                    png_filename = f'image{i}.png'
                    png_rid = 'rId2'
                    svg_rid = 'rId3' if use_compat_mode else 'rId2'

                    shutil.copy(svg_path, media_dir / svg_filename)

                    slide_has_png = False
                    if use_compat_mode:
                        png_path = media_dir / png_filename
                        png_success = convert_svg_to_png(
                            svg_path, png_path,
                            width=pixel_width, height=pixel_height,
                        )
                        if png_success:
                            slide_has_png = True
                            has_any_image = True
                        else:
                            if verbose:
                                print(f"  [{i}/{len(svg_files)}] {svg_path.name} - PNG generation failed, using pure SVG")
                            svg_rid = 'rId2'

                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    slide_xml = create_slide_xml_with_svg(
                        slide_num,
                        png_rid=png_rid, svg_rid=svg_rid,
                        width_emu=width_emu, height_emu=height_emu,
                        transition=transition,
                        transition_duration=transition_duration,
                        auto_advance=auto_advance,
                        use_compat_mode=(use_compat_mode and slide_has_png),
                    )
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'
                    rels_xml = create_slide_rels_xml(
                        png_rid=png_rid, png_filename=png_filename,
                        svg_rid=svg_rid, svg_filename=svg_filename,
                        use_compat_mode=(use_compat_mode and slide_has_png),
                    )
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                # --- Process notes (shared between native and legacy mode) ---
                notes_content = ''
                if enable_notes:
                    svg_stem = svg_path.stem
                    notes_content = notes.get(svg_stem, '') if notes else ''
                    notes_text = markdown_to_plain_text(notes_content) if notes_content else ''
                    if notes_text:
                        notes_slides_dir = extract_dir / 'ppt' / 'notesSlides'
                        notes_slides_dir.mkdir(exist_ok=True)

                        notes_xml_path = notes_slides_dir / f'notesSlide{slide_num}.xml'
                        notes_xml = create_notes_slide_xml(slide_num, notes_text)
                        with open(notes_xml_path, 'w', encoding='utf-8') as f:
                            f.write(notes_xml)

                        notes_rels_dir = notes_slides_dir / '_rels'
                        notes_rels_dir.mkdir(exist_ok=True)
                        notes_rels_path = notes_rels_dir / f'notesSlide{slide_num}.xml.rels'
                        notes_rels_xml = create_notes_slide_rels_xml(slide_num)
                        with open(notes_rels_path, 'w', encoding='utf-8') as f:
                            f.write(notes_rels_xml)

                        _append_relationship(
                            rels_path,
                            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide',
                            f'../notesSlides/notesSlide{slide_num}.xml',
                        )
                        notes_slides_created.add(slide_num)

                if verbose:
                    if use_native_shapes:
                        mode_str = " (Native)"
                    elif use_compat_mode and not use_native_shapes:
                        mode_str = " (PNG+SVG)" if has_any_image else " (SVG)"
                    else:
                        mode_str = " (SVG)"
                    has_notes = slide_num in notes_slides_created
                    notes_str = " +notes" if has_notes else ""
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name}{mode_str}{notes_str}")

                success_count += 1

            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - Error: {e}")

        # Update [Content_Types].xml
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()

        types_to_add: list[str] = []
        if not use_native_shapes:
            if 'Extension="svg"' not in content_types:
                types_to_add.append('  <Default Extension="svg" ContentType="image/svg+xml"/>')
        if has_any_image and 'Extension="png"' not in content_types:
            types_to_add.append('  <Default Extension="png" ContentType="image/png"/>')
        if use_native_shapes and 'Extension="jpg"' not in content_types:
            types_to_add.append('  <Default Extension="jpg" ContentType="image/jpeg"/>')
        if use_native_shapes and 'Extension="jpeg"' not in content_types:
            types_to_add.append('  <Default Extension="jpeg" ContentType="image/jpeg"/>')

        if types_to_add:
            content_types = content_types.replace(
                '</Types>', '\n'.join(types_to_add) + '\n</Types>',
            )
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Add notesSlides content types
        if enable_notes and notes_slides_created:
            for i in sorted(notes_slides_created):
                override = (
                    f'  <Override PartName="/ppt/notesSlides/notesSlide{i}.xml" '
                    f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                )
                if override not in content_types:
                    content_types = content_types.replace('</Types>', override + '\n</Types>')
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Repackage PPTX
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)

        if verbose:
            print()
            print(f"[Done] Saved: {output_path}")
            print(f"  Succeeded: {success_count}, Failed: {len(svg_files) - success_count}")
            if use_compat_mode and has_any_image:
                print(f"  Mode: Office compatibility mode (supports all Office versions)")
                if PNG_RENDERER == 'svglib' and renderer_hint:
                    print(f"  [Tip] {renderer_hint}")

        return success_count == len(svg_files)

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
