#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter

Extracts slide text, tables, speaker notes, and embedded pictures from
Open XML PowerPoint files into Markdown.

Primary use case: PPTX source decks -> Markdown for PPT generation input.

Dependency:
    pip install python-pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SUPPORTED_FORMATS = {
    ".pptx": "PowerPoint Presentation",
    ".pptm": "Macro-enabled PowerPoint Presentation",
    ".ppsx": "PowerPoint Slide Show",
    ".ppsm": "Macro-enabled PowerPoint Slide Show",
    ".potx": "PowerPoint Template",
    ".potm": "Macro-enabled PowerPoint Template",
}


@dataclass
class LeafShape:
    """Flattened leaf shape with stable position ordering."""

    shape: object
    top: int
    left: int


def normalize_text(value: str) -> str:
    """Collapse whitespace while preserving paragraph boundaries elsewhere."""
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"\s+", " ", line).strip() for line in value.split("\n")]
    lines = [line for line in lines if line]
    return "\n".join(lines)


def escape_table_cell(value: str) -> str:
    """Escape Markdown table syntax inside a cell."""
    return normalize_text(value).replace("|", r"\|") or " "


def iter_leaf_shapes(shapes: object) -> list[LeafShape]:
    """Return a flattened, reading-order list of shapes."""
    items: list[LeafShape] = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            items.extend(iter_leaf_shapes(shape.shapes))
            continue
        items.append(
            LeafShape(
                shape=shape,
                top=int(getattr(shape, "top", 0) or 0),
                left=int(getattr(shape, "left", 0) or 0),
            )
        )
    items.sort(key=lambda item: (item.top, item.left))
    return items


def text_frame_to_markdown(text_frame: object) -> str:
    """Convert a PowerPoint text frame into Markdown."""
    paragraphs = []
    visible_paragraphs = [
        paragraph for paragraph in text_frame.paragraphs
        if normalize_text(paragraph.text)
    ]
    if not visible_paragraphs:
        return ""

    list_like = any(paragraph.level > 0 for paragraph in visible_paragraphs)
    if not list_like:
        list_like = len(visible_paragraphs) > 1

    for paragraph in visible_paragraphs:
        text = normalize_text(paragraph.text)
        if not text:
            continue
        if list_like:
            indent = "  " * max(paragraph.level, 0)
            paragraphs.append(f"{indent}- {text}")
        else:
            paragraphs.append(text)

    if list_like:
        return "\n".join(paragraphs)
    return "\n\n".join(paragraphs)


def table_to_markdown(table: object) -> str:
    """Convert a PowerPoint table to a Markdown table."""
    rows = []
    for row in table.rows:
        cells = [escape_table_cell(cell.text) for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    column_count = max(len(row) for row in rows)
    normalized_rows = [row + [" "] * (column_count - len(row)) for row in rows]
    header = normalized_rows[0]
    separator = ["---"] * column_count
    body = normalized_rows[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def save_picture(shape: object, asset_dir: Path, slide_index: int, image_index: int) -> str | None:
    """Persist a picture shape to the output asset directory."""
    try:
        image = shape.image
    except Exception:
        return None

    ext = (image.ext or "png").lower()
    filename = f"slide_{slide_index:02d}_image_{image_index:02d}.{ext}"
    output_path = asset_dir / filename
    output_path.write_bytes(image.blob)
    return filename


def extract_notes(slide: object) -> str:
    """Extract speaker notes text from a slide, if available."""
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""

    blocks = []
    for item in iter_leaf_shapes(notes_slide.shapes):
        shape = item.shape
        if not getattr(shape, "has_text_frame", False):
            continue
        text = text_frame_to_markdown(shape.text_frame)
        if text:
            blocks.append(text)

    return "\n\n".join(blocks).strip()


def convert_presentation_to_markdown(
    input_path: str,
    output_path: str | None = None,
) -> str:
    """Convert a supported PowerPoint file to Markdown."""
    input_file = Path(input_path)
    if not input_file.exists():
        print(f"[ERROR] File not found: {input_path}")
        return ""

    suffix = input_file.suffix.lower()
    if suffix not in SUPPORTED_FORMATS:
        supported = ", ".join(sorted(SUPPORTED_FORMATS.keys()))
        print(f"[ERROR] Unsupported format: {suffix}")
        print(f"   Supported: {supported}")
        print("   Legacy .ppt files should be resaved as .pptx or exported to PDF first.")
        return ""

    print(f"[INFO] Converting {SUPPORTED_FORMATS[suffix]}: {input_file.name}")

    if output_path:
        out_file = Path(output_path)
    else:
        out_file = input_file.with_suffix(".md")

    out_file.parent.mkdir(parents=True, exist_ok=True)
    asset_dir = out_file.parent / f"{out_file.stem}_files"

    presentation = Presentation(str(input_file))
    lines = [
        f"# {input_file.stem}",
        "",
        f"- Source: `{input_file.name}`",
        f"- Total slides: {len(presentation.slides)}",
        "",
    ]

    image_count = 0
    asset_dir_used = False

    for slide_index, slide in enumerate(presentation.slides, 1):
        lines.append(f"## Slide {slide_index}")
        lines.append("")

        blocks = []
        for item in iter_leaf_shapes(slide.shapes):
            shape = item.shape

            if getattr(shape, "has_table", False):
                table_md = table_to_markdown(shape.table)
                if table_md:
                    blocks.append(table_md)
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                next_image_index = image_count + 1
                asset_dir.mkdir(parents=True, exist_ok=True)
                filename = save_picture(shape, asset_dir, slide_index, next_image_index)
                if filename is None:
                    blocks.append(f"> [Image] {getattr(shape, 'name', 'Picture')}")
                    continue

                image_count = next_image_index
                asset_dir_used = True
                blocks.append(f"![Slide {slide_index} Image {image_count}]({asset_dir.name}/{filename})")
                continue

            if getattr(shape, "has_text_frame", False):
                text_md = text_frame_to_markdown(shape.text_frame)
                if text_md:
                    blocks.append(text_md)
                continue

            if getattr(shape, "has_chart", False):
                blocks.append(f"> [Chart] {getattr(shape, 'name', 'Chart')}")

        if blocks:
            lines.append("\n\n".join(blocks))
            lines.append("")
        else:
            lines.append("_No extractable text content._")
            lines.append("")

        notes_md = extract_notes(slide)
        if notes_md:
            lines.append("### Speaker Notes")
            lines.append("")
            lines.append(notes_md)
            lines.append("")

    markdown_content = "\n".join(lines).strip() + "\n"
    out_file.write_text(markdown_content, encoding="utf-8")

    print(f"[OK] Saved Markdown to: {out_file}")
    if asset_dir_used:
        media_files = [path for path in asset_dir.iterdir() if path.is_file()]
        print(f"   Extracted {len(media_files)} image file(s) -> {asset_dir}")

    return markdown_content


def process_directory(input_dir: str, output_dir: str | None = None) -> None:
    """Convert all supported PowerPoint files in a directory to Markdown."""
    input_path = Path(input_dir)

    if output_dir:
        output_path = Path(output_dir)
    else:
        output_path = input_path

    presentation_files = sorted(
        path for path in input_path.iterdir()
        if path.is_file() and path.suffix.lower() in SUPPORTED_FORMATS
    )

    print(f"Found {len(presentation_files)} PowerPoint files")

    for presentation_file in presentation_files:
        output_file = output_path / f"{presentation_file.stem}.md"
        print(f"Processing: {presentation_file.name}")
        result = convert_presentation_to_markdown(str(presentation_file), str(output_file))
        if not result:
            print(f"[WARN] Skipped failed file: {presentation_file.name}")


def main() -> None:
    """Run the CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint files to Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python ppt_to_md.py slides.pptx
  python ppt_to_md.py slides.pptx -o output.md
  python ppt_to_md.py ./decks
  python ppt_to_md.py ./decks -o ./markdown
  python ppt_to_md.py deck.ppsx -o notes/deck.md

Supported formats:
  .pptx  .pptm  .ppsx  .ppsm  .potx  .potm

Legacy .ppt is not parsed directly. Resave it as .pptx or export it to PDF first.
        """,
    )
    parser.add_argument("input", help="Input PowerPoint file or directory")
    parser.add_argument("-o", "--output", help="Output Markdown file or directory path")

    args = parser.parse_args()
    input_path = Path(args.input)

    if input_path.is_file():
        output = args.output or str(input_path.with_suffix(".md"))
        result = convert_presentation_to_markdown(str(input_path), output)
        sys.exit(0 if result else 1)
    if input_path.is_dir():
        process_directory(str(input_path), args.output)
        sys.exit(0)

    print(f"Error: File or directory not found: {args.input}")
    sys.exit(1)


if __name__ == "__main__":
    main()
